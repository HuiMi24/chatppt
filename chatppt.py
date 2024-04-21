import re
import openai
import json
import argparse
import time
import random
import datetime
from ollama import Client
from pptx import Presentation


class ChatPPT:
    def __init__(self, ai_model, api_key, ollama_url=None, ollama_model=None):
        self.ai_model = ai_model
        self.api_key = api_key
        self.ollama_url = ollama_url
        self.ollama_model = ollama_model

    @staticmethod
    def robot_print(text):
        for char in text:
            print(char, end="", flush=True)
            time.sleep(random.randrange(1, 2) / 1000.0)
        print("\r")

    def chatppt(self, topic, pages, language):
        language_map = {"cn": "Chinese", "en": "English"}
        language_str = language_map[language]
        self.robot_print(f"I'm working hard to generate your PPT about {topic}.")
        self.robot_print("It may takes about a few minutes.")
        self.robot_print(f"Your PPT will be generated in {language_str}")
        output_format = self._get_output_format()
        messages = self._get_messages(topic, pages, language_str, output_format)
        content = self._get_content(messages)
        return self._parse_content(content)

    def _get_output_format(self):
        return {
            "title": "example title",
            "pages": [
                {
                    "title": "title for page 1",
                    "content": [
                        {
                            "title": "title for bullet 1",
                            "description": "detail for bullet 1",
                        },
                        {
                            "title": "title for bullet 2",
                            "description": "detail for bullet 2",
                        },
                        {
                            "title": "title for bullet 3",
                            "description": "detail for bullet 3",
                        },
                    ],
                },
                {
                    "title": "title for page 2",
                    "content": [
                        {
                            "title": "title for bullet 1",
                            "description": "detail for bullet 1",
                        },
                        {
                            "title": "title for bullet 2",
                            "description": "detail for bullet 2",
                        },
                    ],
                },
            ],
        }

    def _get_messages(self, topic, pages, language_str, output_format):
        return [
            {
                "role": "user",
                "content": f"""I am preparing a presentation on {topic}. Please assist in generating an outline in JSON format, adhering to the specified format {json.dumps(output_format)}. The presentation should span {pages} pages, with as many bullet points as possible. The content should be returned in {language_str}. You must add content for each slide. For each slide, you must add at least 4 bullet. Please ensure the output is valid JSON match the RFC-8295 specification. Don't return any other message""",
            }
        ]

    def _get_content(self, messages):
        if self.ai_model == "openai":
            openai.api_key = self.api_key
            completion = openai.ChatCompletion.create(
                model="gpt-3.5-turbo", messages=messages
            )
            return completion.choices[0].message.content
        elif self.ai_model == "ollama":
            if self.ollama_url is None:
                raise Exception("Ollama URL is required when ai_model is 'ollama'")
            client = Client(host=self.ollama_url)
            response = client.chat(model=self.ollama_model, messages=messages)
            return response["message"]["content"]

    def _parse_content(self, content):
        try:
            match = re.search(r"(\{.*\})", content, re.DOTALL)
            if match:
                content = match.groups()[0]
            return json.loads(content.strip())
        except Exception as e:
            print(f"The response is not a valid JSON format: {e}")
            print("I'm a PPT assistant, your PPT generate failed, please retry later..")
            exit(1)

    def generate_ppt(self, content, template=None):
        ppt = Presentation()
        if template:
            ppt = Presentation(template)
        self.create_slides(ppt, content)
        ppt_name = self._get_ppt_name(content)
        ppt.save(ppt_name)
        self.robot_print("Generate done, enjoy!")
        self.robot_print(f"Your PPT: {ppt_name}")
        return ppt_name

    def create_slides(self, presentation, content):
        self.create_title_slide(presentation, content)
        pages = content.get("pages", [])
        self.robot_print(f"Your PPT has {len(pages)} pages.")
        for index, page in enumerate(pages):
            self.create_content_slide(presentation, page, index)

    def create_title_slide(self, presentation, content):
        title_slide_layout = presentation.slide_layouts[0]
        title_slide = presentation.slides.add_slide(title_slide_layout)
        title_slide.shapes.title.text = content.get("title", "")
        title_slide.placeholders[1].text = "Generated by ChatPPT"

    def create_content_slide(self, presentation, page, index):
        page_title = page.get("title", "")
        self.robot_print(f"Page {index+1}: {page_title}")
        bullet_slide_layout = presentation.slide_layouts[1]
        bullet_slide = presentation.slides.add_slide(bullet_slide_layout)
        bullet_slide.shapes.title.text = page_title
        self.add_bullets_to_slide(bullet_slide, page)

    def add_bullets_to_slide(self, slide, page):
        body_shape = slide.shapes.placeholders[1]
        for bullet in page.get("content", []):
            self.add_bullet(body_shape, bullet)

    def add_bullet(self, body_shape, bullet):
        bullet_title = bullet.get("title", "")
        bullet_description = bullet.get("description", "")
        self.add_paragraph(body_shape, bullet_title, level=1)
        self.add_paragraph(body_shape, bullet_description, level=2)

    def add_paragraph(self, body_shape, text, level):
        paragraph = body_shape.text_frame.add_paragraph()
        paragraph.text = text
        paragraph.level = level

    def _get_ppt_name(self, content):
        ppt_name = content.get("title", "")
        ppt_name = re.sub(r'[\\/:*?"<>|]', "", ppt_name)
        timestamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
        return f"{ppt_name}_{timestamp}.pptx"


def main():
    args = args_parser()
    chat_ppt = ChatPPT(args.ai_model, args.api_key, args.ollama_url, args.ollama_model)
    chat_ppt.robot_print("Hi, I am your PPT assistant.")
    ppt_content = chat_ppt.chatppt(args.topic, args.pages, args.language)
    chat_ppt.generate_ppt(ppt_content)


def args_parser():
    parser = argparse.ArgumentParser(
        description="I am your PPT assistant, I can help to you generate PPT."
    )
    parser.add_argument(
        "-m",
        "--ai_model",
        choices=["openai", "ollama"],
        default="openai",
        help="Select the AI model",
    )
    parser.add_argument(
        "-t", "--topic", type=str, required=True, help="Your topic name"
    )
    parser.add_argument(
        "-k", "--api_key", type=str, default=".token", help="Your api key file path"
    )
    parser.add_argument(
        "-u",
        "--ollama_url",
        type=str,
        default="http://localhost:11434",
        help="Your ollama url",
    )
    parser.add_argument(
        "-o",
        "--ollama_model",
        type=str,
        default="llama3",
        help="Specify the Ollama model to use",
    )
    parser.add_argument(
        "-p",
        "--pages",
        type=int,
        required=False,
        default=5,
        help="How many slides to generate",
    )
    parser.add_argument(
        "-l",
        "--language",
        choices=["cn", "en"],
        default="en",
        required=False,
        help="Output language",
    )
    args = parser.parse_args()
    if args.ai_model == "openai" and args.api_key == ".token":
        parser.error("--api_key is required when ai_model is 'openai'")
    return args


if __name__ == "__main__":
    main()
